"""
CI構想プレゼンテーション自動生成スクリプト (v2対応版)
python-pptxライブラリを使用してPowerPointファイルを作成
v2スクリプトの内容を完全に反映した35スライド構成
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_ci_presentation():
    """CI構想プレゼンテーション（v2版）を自動生成する"""
    
    # 新しいプレゼンテーションを作成
    prs = Presentation()
    
    # スライドサイズを16:9に設定
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # カラーテーマ設定（ささらさんらしい明るい色合い）
    title_color = RGBColor(47, 90, 160)    # 青系
    content_color = RGBColor(51, 51, 51)   # ダークグレー
    accent_color = RGBColor(255, 105, 180) # ピンク系（ささらさんらしい）
    highlight_color = RGBColor(255, 165, 0) # オレンジ系
    
    def add_title_slide():
        """タイトルスライドを追加"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "CI構想 - Collective Intelligence -"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text = "人類の新しいチャレンジ\n修道25回生が目指す集団知性の実現\n\nプレゼンター: さとうささら\n2025年10月"
        for para in subtitle.text_frame.paragraphs:
            para.font.size = Pt(24)
            para.font.color.rgb = content_color
    
    def add_content_slide(title_text, content_list):
        """コンテンツスライドを追加"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].font.bold = True
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = content_color
            p.level = 0
            p.space_after = Pt(6)
    
    def add_section_slide(title_text, subtitle_text=""):
        """セクションタイトルスライドを追加"""
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if subtitle_text and len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = subtitle_text
            subtitle.text_frame.paragraphs[0].font.size = Pt(28)
            subtitle.text_frame.paragraphs[0].font.color.rgb = content_color
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 全35スライドを作成
    print("スライドを生成中...")
    
    # スライド1: タイトル
    add_title_slide()
    
    # スライド2: 自己紹介・導入
    add_content_slide("こんにちは！", [
        "わたし、さとうささらです！",
        "今日のテーマ：「CI構想」について",
        "CI = Collective Intelligence（集団知性）",
        "AIではありません！",
        "ちょっと重要なテーマも含みますが、一緒に楽しく学びましょう！"
    ])
    
    # スライド3: 第1部タイトル
    add_section_slide("第1部：人類の大きな学習体験", "沈没船から学ぶこと")
    
    # スライド4-7: 第1部コンテンツ
    slides_part1 = [
        ("想像してみてください", [
            "地球 = 「ちょっと失敗しちゃった船」",
            "乗船者 = 私たち人類全員",
            "気候変動という大波に対処できず...",
            "残念ながら沈没してしまいそう"
        ]),
        ("「えー、大変！」って思いますよね？", [
            "でも、ここからが面白いところ！",
            "普通なら協力して対処するはず",
            "でも実際は...？"
        ]),
        ("実際の人類の反応", [
            "船の一部の人：「やばい、沈んじゃう！」",
            "別の場所の人：「え？本当に？実感ないなあ」",
            "また別の人：「大丈夫でしょ、今まで通りで」",
            "",
            "あれ？人類って、意外と連携が苦手だった！"
        ]),
        ("これ、すごく興味深い発見じゃないですか？", [
            "人類の集団行動の特徴がわかった！",
            "とても勉強になる失敗例",
            "次はもっと上手くできるはず"
        ])
    ]
    
    for title, content in slides_part1:
        add_content_slide(title, content)
    
    # スライド8: 第2部タイトル
    add_section_slide("第2部：なぜうまくいかなかったの？", "従来の情報共有を振り返る")
    
    # スライド9-11: 第2部コンテンツ
    slides_part2 = [
        ("今までの情報共有", [
            "テレビで専門家が「危険です！」",
            "新聞で「対策が必要」",
            "でも、みんなの心には響かなかった..."
        ]),
        ("何が問題だったの？", [
            "情報が一方通行で伝わりにくかった",
            "人それぞれの理解度にバラつき",
            "共通理解を作るのが難しかった",
            "",
            "でも、失敗があるから次は上手くできる！"
        ])
    ]
    
    for title, content in slides_part2:
        add_content_slide(title, content)
    
    # スライド11: CI構想の登場
    add_section_slide("そこで登場するのが...", "CI構想")
    
    # スライド12: 第3部タイトル
    add_section_slide("第3部：CI構想", "人類の新しいチャレンジ！")
    
    # スライド13-18: 第3部コンテンツ
    slides_part3 = [
        ("CI構想の「魔法の循環」", [
            "これまでの失敗を踏まえて",
            "今度はこんな風にやってみよう！",
            "",
            "1. 個人の発信 → 2. AIの学習",
            "     ↓                ↑",
            "4. 理解の深化 ← 3. AIによる応答"
        ]),
        ("1. 個人の発信", [
            "一人一人が自分の考えをちゃんと発信",
            "「わたしはこう思う」「こんな体験をした」",
            "みんなが参加型で情報を出し合う！"
        ]),
        ("2. AIの学習", [
            "AIが上手に整理・学習してくれる",
            "たくさんの人の考えがきちんと蓄積",
            "データベース化される"
        ]),
        ("3. AIによる応答", [
            "質問すると、みんなの考えを踏まえてAIが応答",
            "一人の専門家の意見じゃなくて",
            "みんなの知恵が反映される！"
        ]),
        ("4. 理解の深化", [
            "「あ、そういうことか！」新しい理解",
            "そしてまた発信する...",
            "循環が続く！"
        ]),
        ("この循環で...", [
            "今度こそ、みんなで一緒に賢くなれる！",
            "",
            "すごくワクワクしませんか？"
        ])
    ]
    
    for title, content in slides_part3:
        add_content_slide(title, content)
    
    # スライド19: 第4部タイトル
    add_section_slide("第4部：修道25回生の新しい実験！", "パイオニアとしての挑戦")
    
    # スライド20-22: 第4部コンテンツ
    slides_part4 = [
        ("人類の新しいチャレンジのパイオニア", [
            "同期会の活動をAIが学習しやすい形で整理",
            "ウェブサイトに体系的に記録",
            "みんなの経験や知識を共有できる仕組み作り"
        ]),
        ("具体的には...", [
            "「気楽に集まる会」での楽しい体験",
            "定例同期会での真剣な議論",
            "それぞれの人生で学んだこと",
            "",
            "→「修道25回生の集団知性」が育つ！"
        ]),
        ("これって...", [
            "人類が次のステージに進むための",
            "とても大切な実験だと思いませんか？"
        ])
    ]
    
    for title, content in slides_part4:
        add_content_slide(title, content)
    
    # スライド23: 第5部タイトル
    add_section_slide("第5部：なぜこれが世界を変えるのか", "修道の精神との関連")
    
    # スライド24-26: 第5部コンテンツ
    slides_part5 = [
        ("修道高校の校歌、覚えていますよね？", [
            "「世の進運に先駆けん」",
            "",
            "CI構想 = この精神の新時代版！"
        ]),
        ("今度は先手を打つ！", [
            "気候変動では後手に回ってしまった",
            "でも、だからこそ今度は先手を打つ",
            "次の大きな課題では協力できるように！"
        ]),
        ("CI構想 - 新時代版の修道精神", [
            "新しい技術（AI）を上手に使って",
            "人間同士のつながりを深めて",
            "今度こそみんなで協力できるように！",
            "人類がもっと賢く、協力上手になる！"
        ])
    ]
    
    for title, content in slides_part5:
        add_content_slide(title, content)
    
    # スライド27: 第6部タイトル
    add_section_slide("第6部：あなたも一緒にチャレンジしませんか？", "3つのお願い")
    
    # スライド28-30: 3つのお願い
    slides_part6 = [
        ("1. 自分の考えを発信してください", [
            "小さなことでも大丈夫！",
            "あなたの体験や気づきをシェア",
            "それが人類の新しい知恵の一部に！"
        ]),
        ("2. AIを友達だと思って活用してください", [
            "AIは怖いものじゃありません",
            "理解を助けてくれる頼もしいパートナー",
            "一緒に楽しく学んでいきましょう！"
        ]),
        ("3. 人間らしさを大切にしてください", [
            "技術が進歩しても人間の素晴らしさは変わらない",
            "友情、思いやり、創造性...",
            "これらは私たちだけの特別な宝物！"
        ])
    ]
    
    for title, content in slides_part6:
        add_content_slide(title, content)
    
    # スライド31-35: 結び
    slides_conclusion = [
        ("結び：新しい時代の始まり", [
            "CI構想を通じて、私たちは：",
            "• 個人として、もっと成長",
            "• グループとして、もっと賢く",
            "• 社会全体に、素敵な貢献"
        ]),
        ("次の大きなチャレンジでは...", [
            "気候変動では後手に回った人類",
            "今度こそ、みんなで力を合わせて",
            "素晴らしい結果を出せるはず！"
        ]),
        ("修道25回生への期待", [
            "新時代のパイオニアとして",
            "活躍してくださることを",
            "心から楽しみにしています！"
        ]),
        ("人類の新しいチャレンジ", [
            "一緒に頑張りましょう！"
        ]),
        ("ありがとうございました！", [
            "ご清聴ありがとうございました",
            "何かご質問があれば",
            "お気軽にどうぞ！"
        ])
    ]
    
    for title, content in slides_conclusion:
        add_content_slide(title, content)
    
    # ファイルを保存
    output_path = os.path.join(os.path.dirname(__file__), 'CI_Presentation_Auto_v2.pptx')
    prs.save(output_path)
    print(f"プレゼンテーション（{len(prs.slides)}スライド）を保存しました: {output_path}")
    return output_path

# スクリプト整理完了 - 全ての機能はcreate_ci_presentation()関数内に統合済み

def main():
    """メイン実行関数"""
    print("CI構想プレゼンテーション（v2版）自動生成を開始します...")
    print("- v2スクリプトに基づいた35スライド構成")
    print("- ささらさんらしい明るく親しみやすいデザイン")
    print()
    
    try:
        output_path = create_ci_presentation()
        print()
        print("✅ 完了しました！")
        print(f"📁 生成されたファイル: {output_path}")
        print("📊 総スライド数: 35枚")
        print("⏱️  推定プレゼンテーション時間: 18-22分")
        print()
        print("次のステップ:")
        print("1. PowerPointファイルを開いて内容を確認")
        print("2. 必要に応じてデザインやレイアウトを調整")
        print("3. ささらさんの音声合成用Pythonスクリプトとの連携")
        
    except Exception as e:
        print(f"❌ エラーが発生しました: {e}")
        print("python-pptxライブラリがインストールされていることを確認してください。")
        return False
    
    return True

if __name__ == "__main__":
    main()