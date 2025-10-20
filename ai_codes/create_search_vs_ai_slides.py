"""
検索とAI応答の違いを説明するPowerPoint自動生成スクリプト
python-pptxライブラリを使用
最小3枚構成
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_search_vs_ai_slides():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title_color = RGBColor(47, 90, 160)
    content_color = RGBColor(51, 51, 51)
    accent_color = RGBColor(255, 105, 180)

    def add_content_slide(title_text, content_list, note=None):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].font.bold = True
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        for i, item in enumerate(content_list):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = content_color
            p.level = 0
            p.space_after = Pt(6)
        if note:
            p = text_frame.add_paragraph()
            p.text = f"【図案メモ】{note}"
            p.font.size = Pt(14)
            p.font.color.rgb = accent_color
            p.level = 0

    # スライド1（供給側視点の違い）
    add_content_slide(
        "情報の集め方の違い（供給側視点）",
        [
            "Google検索: クローラーがウェブ全体を定期的に巡回・収集し、インデックスを更新",
            "AI（ChatGPTなど）: 過去のウェブ情報を一括で学習し、モデル化。新しい情報は再学習まで反映されない"
        ],
        note="クローラーが常時巡回 vs AIが一括学習して止まる"
    )

    # スライド2（新しさ・反映速度）
    add_content_slide(
        "情報の新しさ・反映速度",
        [
            "Google検索: 新情報は数時間～数日でインデックスに反映",
            "AI: 学習時点で知識が止まる（半年～1年前の情報が多い）"
        ],
        note="タイムライン上に『検索は最新』『AIは過去でストップ』の図"
    )

    # スライド3（CI構想の情報循環）
    add_content_slide(
        "AIにおける情報循環（CI構想の視点）",
        [
            "（個人の）ウェブページへの発信 → AIの整理・学習 → AIによる応答 → （個人の）理解への影響",
            "この循環を通じて、個人と集団の知性が高まる"
        ],
        note="情報循環の流れ図（矢印でつなぐ）"
    )

    output_path = os.path.join(os.path.dirname(__file__), 'Search_vs_AI_Slides.pptx')
    prs.save(output_path)
    print(f"検索とAI応答の違いスライド（3枚）を保存しました: {output_path}")
    return output_path

if __name__ == "__main__":
    create_search_vs_ai_slides()
