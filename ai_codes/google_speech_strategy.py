from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_google_speech_presentation():
    # プレゼンテーションオブジェクトを作成
    prs = Presentation()
    
    # スライド1: タイトルスライド
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Googleが音声認識を無料提供する戦略的メリット"
    subtitle.text = "Google Speech Recognition APIの無料戦略を分析"
    
    # スライド2: 概要
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "概要"
    tf = content.text_frame
    tf.text = "Googleは音声認識APIを無料で提供"
    
    p = tf.add_paragraph()
    p.text = "一見すると収益に直結しない施策"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "実は長期的な戦略投資"
    p.level = 0
    run = p.runs[0]
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 0, 0)
    
    # スライド3: 6つの戦略的メリット
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "6つの戦略的メリット"
    tf = content.text_frame
    
    benefits = [
        "データ収集とAI改善",
        "有料版への誘導（フリーミアム）",
        "エコシステムの拡大",
        "技術力のアピール",
        "開発者コミュニティの育成",
        "市場シェア確保"
    ]
    
    tf.text = f"1. {benefits[0]}"
    for i, benefit in enumerate(benefits[1:], 2):
        p = tf.add_paragraph()
        p.text = f"{i}. {benefit}"
        p.level = 0
    
    # スライド4-9: 各メリットの詳細
    details = [
        {
            "title": "1. データ収集とAI改善",
            "points": [
                "世界中のユーザーの音声データを収集",
                "多様な言語・方言・アクセントの蓄積",
                "機械学習モデルの訓練データとして活用",
                "→ 音声認識精度の継続的向上"
            ]
        },
        {
            "title": "2. 有料版への誘導（フリーミアム）",
            "points": [
                "無料版で体験してもらう",
                "商用・大規模利用時に有料版へアップグレード",
                "Google Cloud Speech APIの売上拡大",
                "典型的なフリーミアム戦略"
            ]
        },
        {
            "title": "3. エコシステムの拡大",
            "points": [
                "開発者が気軽に音声機能を実装",
                "Googleのテクノロジーに慣れ親しんでもらう",
                "Google Cloudの他サービス利用促進",
                "技術的囲い込み効果"
            ]
        },
        {
            "title": "4. 技術力のアピール",
            "points": [
                "Googleの音声認識技術の優秀さを実証",
                "競合他社への技術的優位性の示威",
                "ブランド価値の向上",
                "技術リーダーシップの確立"
            ]
        },
        {
            "title": "5. 開発者コミュニティの育成",
            "points": [
                "音声技術を使ったイノベーション促進",
                "新しいサービス・アプリの創出",
                "技術的パートナーの獲得",
                "長期的な技術エコシステム構築"
            ]
        },
        {
            "title": "6. 市場シェア確保",
            "points": [
                "音声認識市場でのデファクトスタンダード化",
                "他社技術への流出防止",
                "市場における優位性の確保",
                "将来的な収益基盤の構築"
            ]
        }
    ]
    
    for detail in details:
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = detail["title"]
        tf = content.text_frame
        tf.text = detail["points"][0]
        
        for point in detail["points"][1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            if point.startswith("→"):
                run = p.runs[0]
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 128, 0)
    
    # スライド10: まとめ
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "まとめ"
    tf = content.text_frame
    tf.text = "短期的収益 < 長期的戦略投資"
    # 最初の段落のrun経由でフォーマットを変更
    first_paragraph = tf.paragraphs[0]
    run = first_paragraph.runs[0]
    run.font.bold = True
    run.font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "無料提供により将来的な市場支配を狙う"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "デベロッパー・エコシステムの構築が鍵"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "音声認識分野での覇権確立戦略"
    p.level = 0
    run = p.runs[0]
    run.font.color.rgb = RGBColor(255, 0, 0)
    
    # プレゼンテーションを保存
    prs.save('google_speech_strategy.pptx')
    print("PowerPointファイル 'google_speech_strategy.pptx' を作成しました。")

if __name__ == "__main__":
    # python-pptxが必要です: pip install python-pptx
    create_google_speech_presentation()