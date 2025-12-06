"""
CI構想プレゼンテーション自動生成スクリプト (v3最終版)
presentation-structure-v3.mdに基づいた15分版
python-pptxライブラリを使用してPowerPointファイルを作成

特徴：
- 15分プレゼンテーション用（12-14スライド）
- 2つの軸で情報環境の変化を整理
- エネルギー比較でLLMの「理解」を説明
- 問いを共有するトーン
- 簡潔で要点を絞った構成
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_ci_presentation_v3_final():
    """CI構想プレゼンテーション（v3最終版）を自動生成する"""
    
    # 新しいプレゼンテーションを作成
    prs = Presentation()
    
    # スライドサイズを16:9に設定
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # カラーテーマ設定
    title_color = RGBColor(47, 90, 160)    # 青系
    content_color = RGBColor(51, 51, 51)   # ダークグレー
    accent_color = RGBColor(220, 50, 50)   # 赤系（強調用）
    
    def add_title_slide():
        """タイトルスライドを追加"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "CI構想\nCollective Intelligence"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text = "集団の知性を高めるための\n新しいチャレンジ\n\n修道25回生 同期会\n2025年11月"
        for para in subtitle.text_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = content_color
    
    def add_content_slide(title_text, content_list, note=""):
        """コンテンツスライドを追加"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].font.bold = True
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # レベル判定（インデント制御）
            level = 0
            text = item
            if item.startswith("  - "):
                level = 1
                text = item[4:]
            elif item.startswith("- "):
                level = 0
                text = item[2:]
            
            p.text = text
            p.font.size = Pt(18)
            p.font.color.rgb = content_color
            p.level = level
            p.space_after = Pt(8)
        
        # ノートを追加
        if note:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = note
    
    def add_section_slide(title_text, subtitle_text=""):
        """セクションタイトルスライドを追加"""
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if subtitle_text and len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = subtitle_text
            subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle.text_frame.paragraphs[0].font.color.rgb = content_color
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # スライド作成開始
    print("スライドを生成中...")
    
    # スライド1: タイトル
    add_title_slide()
    
    # スライド2: 本日の目的
    add_content_slide("本日の目的", [
        "「考え始めるきっかけ」を与える",
        "",
        "- 答えを提示するのではなく、問いを共有する",
        "- 説得するのではなく、一緒に考える仲間を増やす",
        "- 完璧な構想ではなく、進行中の試みとして提示"
    ], note="視覚要素: 対話のイメージ、考える人々")
    
    # スライド3: 2つの軸（第1の軸）
    add_content_slide("第1の軸：情報伝達の効率化・理解の深さ", [
        "我々70歳は、3つの時代を体験した唯一の世代",
        "",
        "- インターネット以前：「見る、聞く、読む」（受動的）",
        "- インターネット検索：「探す、選ぶ」（能動的）",
        "- AI対話：「対話し、理解を深める」（対話的）",
        "",
        "重要な気づき：「見る、聞く」だけでは正確な理解はできない"
    ], note="視覚要素: 3つの時代を示す表、矢印で進化を表現")
    
    # スライド4: 2つの軸（第2の軸）
    add_content_slide("第2の軸：情報発信の多様化", [
        "インターネットがもたらした変化",
        "",
        "- インターネット以前：マスメディアのみ",
        "- インターネット時代：誰もが発信可能",
        "  - YouTuber、個人ブログ、SNS",
        "- 現在の象徴例：NHK ONE（2025年10月～）",
        "",
        "新たな課題：情報の質のばらつき、選択の難しさ"
    ], note="視覚要素: マスメディア→個人発信の変化を図示")
    
    # スライド5: 問いかけ - 気候変動
    add_section_slide("問いかけ", "なぜ気候変動対策は手遅れになったのか？")
    
    # スライド6: 気候変動の失敗から学ぶ
    add_content_slide("気候変動対策の失敗から学ぶ", [
        "沈没寸前の船の比喩",
        "",
        "- マスメディア時代 = 「見る、聞く」だけの限界",
        "  - 情報は一方向、受動的な理解では不十分",
        "",
        "- 「全員が正しく理解する」には何が必要だったか？",
        "  - 「探す、選ぶ」（能動性）",
        "  - さらには「対話し、理解を深める」（深化）",
        "",
        "問い：同じ過ちを繰り返さないために、何が必要か？"
    ], note="視覚要素: 船のイラスト、一方向の矢印")
    
    # スライド7: Google検索とAI対話の違い（仕組み）
    add_content_slide("Google検索とAI対話の違い（仕組み）", [
        "Google検索",
        "- インデックス作成（キーワード対応表）",
        "- ページのリストを提示",
        "- 体験：「探す、選ぶ」",
        "",
        "LLM（ChatGPTなど）",
        "- 言語理解の学習（意味・文脈・関係性）",
        "- 対話的な回答を生成",
        "- 体験：「対話し、理解を深める」"
    ], note="視覚要素: 2カラム比較表")
    
    # スライド8: エネルギーで見る「理解」のコスト
    add_content_slide("エネルギーで見る「理解」のコスト", [
        "なぜLLMはこんなに電力を使うのか？",
        "",
        "- 1回の利用：Google検索 0.0003 kWh / LLM 0.003-0.03 kWh",
        "- 理解は探すことの10-100倍のエネルギーが必要",
        "",
        "それは「理解」しているから",
        "- Google検索：「どこにあるか」を教える（図書館の目録）",
        "- LLM：「何を意味するか」を理解して説明する（教師）",
        "",
        "人間の脳も同じ：深く理解するほど多くのエネルギーを使う"
    ], note="視覚要素: エネルギー比較のグラフ、電球のイラスト")
    
    # スライド9: AI対話時代とCI構想
    add_section_slide("ヒント", "AI対話時代とCI構想")
    
    # スライド10: CI構想の核心
    add_content_slide("2つの軸の統合的発展 = CI構想", [
        "次に来るのは「対話し、理解を深める」時代",
        "",
        "- 第1の軸：AI対話による理解の深化",
        "- 第2の軸：個人発信の質の向上（AIの支援も受けて）",
        "- → 両軸が統合されたCI構想",
        "",
        "CI（Collective Intelligence）：集合知",
        "- 個人発信 → AI学習 → 他者への還元という循環",
        "- 深い理解を共有することで集団の知性を高める"
    ], note="視覚要素: 循環を示す矢印、2つの軸の統合図")
    
    # スライド11: 同期会での実践例
    add_content_slide("同期会での小さな実践例", [
        "活動をウェブで発信",
        "- 第2の軸：発信の多様化を実践",
        "",
        "AIが学習しやすい形で整理",
        "- 構造化された情報として蓄積",
        "",
        "AI対話を通じて次世代へ継承",
        "- 第1の軸：理解の深化を促進",
        "",
        "問い：我々にできることは何か？"
    ], note="視覚要素: Webサイトのスクリーンショット、循環図")
    
    # スライド12: まとめ - 2つの軸の再確認
    add_content_slide("まとめ：2つの軸で見た情報環境の進化", [
        "第1の軸：理解の深さ",
        "- 「見る、聞く」→「探す、選ぶ」→「対話し、理解を深める」",
        "- より深い理解とエネルギーが必要に",
        "",
        "第2の軸：発信の多様化",
        "- マスメディアの独占 → 誰もが発信者に",
        "",
        "2つの軸の統合",
        "- AI対話時代は、両軸が統合される時代",
        "- 深い理解と多様な発信が循環するCI構想"
    ], note="視覚要素: 2つの軸の統合図、上昇矢印")
    
    # スライド13: 我々の世代の役割
    add_content_slide("我々の世代の役割", [
        "3つの時代を体験できる唯一の世代",
        "2つの軸の変化を体験した世代",
        "",
        "この視点を次世代に伝える責任",
        "次の大きな波はすぐそこに",
        "",
        "行動への呼びかけ",
        "- 答えは一つではない",
        "- 大事なのは「考え続けること」と「小さく始めること」",
        "- 今日の話をきっかけに、それぞれが考え、行動してほしい"
    ], note="視覚要素: 世代のイメージ、前進する人々")
    
    # スライド14: 終了
    add_content_slide("ご清聴ありがとうございました", [
        "CI構想：進行中の試み",
        "",
        "一緒に考え、行動しましょう",
        "",
        "ご質問・ご意見をお待ちしています"
    ], note="視覚要素: 感謝のメッセージ、対話のイメージ")
    
    # ファイルを保存
    output_path = os.path.join(os.path.dirname(__file__), 'CI_Presentation_v3_Final.pptx')
    prs.save(output_path)
    print(f"プレゼンテーション（{len(prs.slides)}スライド）を保存しました: {output_path}")
    return output_path

def main():
    """メイン実行関数"""
    print("=" * 60)
    print("CI構想プレゼンテーション（v3最終版）自動生成")
    print("=" * 60)
    print()
    print("【構成の特徴】")
    print("- presentation-structure-v3.mdに完全準拠")
    print("- 15分プレゼンテーション用（14スライド）")
    print("- 2つの軸で情報環境の変化を整理")
    print("- エネルギー比較でLLMの「理解」を説明")
    print("- 問いを共有するトーン")
    print()
    
    try:
        output_path = create_ci_presentation_v3_final()
        print()
        print("=" * 60)
        print("✅ 完了しました！")
        print("=" * 60)
        print(f"📁 生成されたファイル: {output_path}")
        print(f"📊 総スライド数: 14枚")
        print(f"⏱️  推定プレゼンテーション時間: 15分")
        print()
        print("【スライド構成】")
        print("1. タイトル")
        print("2. 本日の目的")
        print("3-4. 2つの軸（第1の軸・第2の軸）")
        print("5-6. 問いかけ：気候変動対策の失敗")
        print("7-8. Google検索とAI対話の違い")
        print("9-10. AI対話時代とCI構想")
        print("11. 同期会での実践例")
        print("12. まとめ：2つの軸の再確認")
        print("13. 我々の世代の役割")
        print("14. 終了")
        print()
        print("【次のステップ】")
        print("1. PowerPointファイルを開いて内容を確認")
        print("2. 必要に応じてデザインやレイアウトを調整")
        print("3. ノート欄の視覚要素指示を参考に画像や図表を追加")
        print("4. タイミングを確認しながらリハーサル")
        print()
        
    except Exception as e:
        print()
        print("=" * 60)
        print("❌ エラーが発生しました")
        print("=" * 60)
        print(f"エラー内容: {e}")
        print()
        print("【対処方法】")
        print("python-pptxライブラリがインストールされていることを確認してください。")
        print()
        print("インストールコマンド:")
        print("  pip install python-pptx")
        print()
        return False
    
    return True

if __name__ == "__main__":
    main()
