"""
CI構想プレゼンテーション自動生成スクリプト (v3対応版)
python-pptxライブラリを使用してPowerPointファイルを作成
v3スクリプトの内容を完全に反映した25スライド構成

改訂のポイント：
- スライド数を35枚→25枚に削減
- 船の比喩表現を軽く調整
- AI vs. CIの説明を明確化
- 具体例を追加
- ささらの天然キャラを強化
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_ci_presentation_v3():
    """CI構想プレゼンテーション（v3版）を自動生成する"""
    
    # 新しいプレゼンテーションを作成
    prs = Presentation()
    
    # スライドサイズを16:9に設定
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # カラーテーマ設定（ささらさんらしい明るい色合い）
    title_color = RGBColor(47, 90, 160)    # 青系
    content_color = RGBColor(51, 51, 51)   # ダークグレー
    accent_color = RGBColor(255, 105, 180) # ピンク系（ささらさんらしい）
    highlight_color = RGBColor(255, 165, 0) # オレンジ系
    
    def add_title_slide():
        """タイトルスライドを追加"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "CI構想 - Collective Intelligence -"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text = "人類の新しいチャレンジ\n修道25回生が目指す集団知性の実現\n\nプレゼンター: さとうささら\n2025年10月"
        for para in subtitle.text_frame.paragraphs:
            para.font.size = Pt(24)
            para.font.color.rgb = content_color
    
    def add_content_slide(title_text, content_list, note=""):
        """コンテンツスライドを追加"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].font.bold = True
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = content_color
            p.level = 0
            p.space_after = Pt(6)
        
        # ノートを追加（視覚要素の指示など）
        if note:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = note
    
    def add_section_slide(title_text, subtitle_text="", note=""):
        """セクションタイトルスライドを追加"""
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if subtitle_text and len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = subtitle_text
            subtitle.text_frame.paragraphs[0].font.size = Pt(28)
            subtitle.text_frame.paragraphs[0].font.color.rgb = content_color
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # ノートを追加
        if note:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = note
    
    # 全25スライドを作成
    print("スライドを生成中...")
    
    # スライド1: タイトル
    add_title_slide()
    
    # スライド2: 自己紹介・導入
    add_content_slide("こんにちは！", [
        "わたし、さとうささらです！",
        "今日のテーマ：「CI構想」について",
        "CI = Collective Intelligence（集団知性）",
        "AIを使うけど、AI「だけ」じゃありません！",
        "人間＋AI＝もっとすごい！",
        "一緒に楽しく学びましょう！"
    ], note="視覚要素: ささらの笑顔、「人間＋AI」の図解")
    
    # スライド3: 第1部タイトル
    add_section_slide("第1部：人類の大きな学習体験", "船の話から学ぶこと",
                     note="視覚要素: 船のシルエット、波のイラスト")
    
    # スライド4-6: 第1部コンテンツ
    add_content_slide("想像してみてください", [
        "地球 = 「今、ちょっとピンチな船」",
        "乗船者 = 私たち人類全員",
        "気候変動という大きな波がきちゃって...",
        "みんなで協力すれば何とかなるはずなのに",
        "うまく連携できなかった"
    ], note="視覚要素: 船のイラスト、波、人々のシルエット")
    
    add_content_slide("実際の人類の反応は...？", [
        "船の一部の人：「やばい、大変！」",
        "別の場所の人：「え？本当？実感ないなあ」",
        "また別の人：「大丈夫でしょ、今まで通りで」",
        "",
        "あれれ？なんか変だな...って思いませんか？（笑）",
        "人類って、意外と連携が苦手だったんです！"
    ], note="視覚要素: 3つの吹き出しイラスト、困った顔のささら")
    
    add_content_slide("でも、これってすごく勉強になる！", [
        "人類の集団行動の特徴がわかった！",
        "とても貴重な学習体験",
        "次はもっと上手くできるはず",
        "ここからが面白いんです！"
    ], note="視覚要素: 電球マーク、前向きな雰囲気")
    
    # スライド7: 第2部タイトル
    add_section_slide("第2部：なぜうまくいかなかったの？", "従来の情報共有を振り返る",
                     note="視覚要素: 疑問符、考えるポーズのささら")
    
    # スライド8: 第2部コンテンツ
    add_content_slide("今までの情報共有の問題点", [
        "【従来の方法】",
        "・テレビで専門家が「危険です！」",
        "・新聞で「対策が必要」",
        "・でも、みんなの心には響かなかった...",
        "",
        "【何が問題？】",
        "・情報が一方通行",
        "・人それぞれの理解度にバラつき",
        "・共通理解を作るのが難しかった",
        "・でも、失敗から学べます！"
    ], note="視覚要素: 左右2カラム、矢印で一方通行を表現")
    
    # スライド9: CI構想の登場
    add_section_slide("そこで登場！", "CI構想\n人類の新しいチャレンジ！",
                     note="視覚要素: キラキラエフェクト、ワクワク感")
    
    # スライド10: 第3部タイトル
    add_section_slide("第3部：CI構想の「魔法の循環」", "これまでの失敗を踏まえて",
                     note="視覚要素: 循環を示す円形の矢印")
    
    # スライド11: 循環図
    add_content_slide("CI構想の4つのステップ", [
        "1. 個人の発信 →",
        "2. AIの学習 →",
        "3. AIによる応答 →",
        "4. 理解の深化 → （1に戻る）",
        "",
        "この循環で、みんなで一緒に賢くなれる！"
    ], note="視覚要素: 4つのステップを示す循環図、アニメーション効果")
    
    # スライド12-14: ステップ詳細
    add_content_slide("1. 個人の発信", [
        "一人一人が自分の考えをちゃんと発信",
        "「わたしはこう思う」",
        "「こんな体験をした」",
        "みんなが参加型で情報を出し合う！"
    ], note="視覚要素: 人々が発信する様子、吹き出し")
    
    add_content_slide("2. AIの学習", [
        "AIが上手に整理・学習してくれる",
        "AIって本当に賢いですよね！",
        "わたしより整理上手かも（笑）",
        "たくさんの人の考えがきちんと蓄積"
    ], note="視覚要素: AIのアイコン、データベースのイメージ")
    
    add_content_slide("3. AIによる応答 → 4. 理解の深化", [
        "【ステップ3】",
        "・質問すると、みんなの考えを踏まえてAIが応答",
        "・一人の専門家じゃなく、みんなの知恵が反映！",
        "",
        "【ステップ4】",
        "・「あ、そういうことか！」新しい理解",
        "・そしてまた発信する...循環が続く！"
    ], note="視覚要素: 質問→応答→理解のフロー図")
    
    # スライド15: ワクワクする結果
    add_content_slide("この循環で...", [
        "今度こそ、みんなで一緒に賢くなれる！",
        "すごくワクワクしませんか？",
        "これが「集団知性」の力です！"
    ], note="視覚要素: 明るい色彩、上昇矢印、笑顔")
    
    # スライド16: 第4部タイトル
    add_section_slide("第4部：修道25回生の実験！", "パイオニアとしての挑戦",
                     note="視覚要素: 修道のロゴ、パイオニアのイメージ")
    
    # スライド17-18: 第4部コンテンツ
    add_content_slide("私たちの具体的な活動", [
        "・同期会の活動をAIが学習しやすい形で整理",
        "・ウェブサイトに体系的に記録",
        "・みんなの経験や知識を共有できる仕組み",
        "",
        "【具体例】",
        "・「気楽に集まる会」→ 競馬を楽しむ会の報告",
        "・定例同期会 → 真剣な議論の記録",
        "・人生経験 → それぞれの知恵の共有"
    ], note="視覚要素: 実際のウェブサイトのスクリーンショット、活動写真")
    
    add_content_slide("「修道25回生の集団知性」が育つ！", [
        "一人一人の経験が集まって",
        "みんなで学び合って",
        "人類の新しい実験",
        "",
        "とても大切なチャレンジだと思いませんか？"
    ], note="視覚要素: 木が育つイメージ、つながる人々")
    
    # スライド19: 第5部タイトル
    add_section_slide("第5部：修道精神とCI構想", "「世の進運に先駆けん」の新しい形",
                     note="視覚要素: 修道の校章、未来へ向かう矢印")
    
    # スライド20: 第5部コンテンツ
    add_content_slide("校歌、覚えていますよね？", [
        "「世の進運に先駆けん」",
        "気候変動では後手に回ってしまった",
        "でも、だからこそ今度は先手を打つ！",
        "次の大きな課題では協力できるように",
        "",
        "【CI構想 = 修道精神の新時代版】",
        "・新しい技術（AI）を上手に使って",
        "・人間同士のつながりを深めて",
        "・人類がもっと賢く、協力上手になる！"
    ], note="視覚要素: 校歌の一節、未来へのイメージ")
    
    # スライド21: 第6部タイトル
    add_section_slide("第6部：あなたも一緒にチャレンジしませんか？", "3つのお願い",
                     note="視覚要素: 招待する手のジェスチャー")
    
    # スライド22: 3つのお願い
    add_content_slide("3つのお願い", [
        "【1. 自分の考えを発信してください】",
        "・小さなことでも大丈夫！",
        "・あなたの体験や気づきをシェア",
        "",
        "【2. AIを友達だと思って活用してください】",
        "・AIは怖いものじゃありません",
        "・理解を助けてくれる頼もしいパートナー",
        "",
        "【3. 人間らしさを大切にしてください】",
        "・友情、思いやり、創造性",
        "・これらは私たちだけの特別な宝物！"
    ], note="視覚要素: 3つの項目をアイコンで表示")
    
    # スライド23: 結び - CI構想のまとめ
    add_content_slide("CI構想を通じて", [
        "・個人として、もっと成長",
        "・グループとして、もっと賢く",
        "・社会全体に、素敵な貢献",
        "",
        "気候変動では後手に回った人類",
        "でも今度こそ、次のチャレンジでは",
        "みんなで力を合わせて素晴らしい結果を！"
    ], note="視覚要素: 成長曲線、明るい未来のイメージ")
    
    # スライド24: 修道25回生への期待
    add_content_slide("修道25回生のパイオニアとして", [
        "新時代のパイオニアとして",
        "人類の新しいチャレンジに挑戦",
        "",
        "一緒に頑張りましょう！",
        "",
        "活躍を心から楽しみにしています！"
    ], note="視覚要素: 集合写真イメージ、前進する人々")
    
    # スライド25: 終了・質疑応答
    add_content_slide("ありがとうございました！", [
        "ご清聴ありがとうございました",
        "何かご質問があれば",
        "お気軽にどうぞ！",
        "",
        "わたし、さとうささらでした！"
    ], note="視覚要素: ささらの笑顔、拍手のイラスト")
    
    # ファイルを保存
    output_path = os.path.join(os.path.dirname(__file__), 'CI_Presentation_Auto_v3.pptx')
    prs.save(output_path)
    print(f"プレゼンテーション（{len(prs.slides)}スライド）を保存しました: {output_path}")
    return output_path

def main():
    """メイン実行関数"""
    print("CI構想プレゼンテーション（v3版）自動生成を開始します...")
    print("- v3スクリプトに基づいた25スライド構成")
    print("- スライド数を35枚→25枚に削減")
    print("- 船の表現を軽く調整")
    print("- AI vs. CIの説明を明確化")
    print("- 具体例（競馬会など）を追加")
    print("- ささらの天然キャラを強化")
    print()
    
    try:
        output_path = create_ci_presentation_v3()
        print()
        print("✅ 完了しました！")
        print(f"📁 生成されたファイル: {output_path}")
        print("📊 総スライド数: 25枚")
        print("⏱️  推定プレゼンテーション時間: 18-22分")
        print()
        print("主な改善点:")
        print("1. スライド数削減（35→25枚）でテンポアップ")
        print("2. 船の比喩を「沈没」→「ピンチ」に軽減")
        print("3. 「人間＋AI＝もっとすごい」を明示")
        print("4. 競馬を楽しむ会などの具体例を追加")
        print("5. ささらの天然キャラを強化")
        print()
        print("次のステップ:")
        print("1. PowerPointファイルを開いて内容を確認")
        print("2. 必要に応じてデザインやレイアウトを調整")
        print("3. ノート欄の視覚要素指示を参考に画像を追加")
        
    except Exception as e:
        print(f"❌ エラーが発生しました: {e}")
        print("python-pptxライブラリがインストールされていることを確認してください。")
        print("インストール: pip install python-pptx")
        return False
    
    return True

if __name__ == "__main__":
    main()
